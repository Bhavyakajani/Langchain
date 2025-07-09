from langchain_community.document_loaders import PyPDFLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from pptx import Presentation
import cv2
import pytesseract

# document_processing.py

def load_file(path):
    """Load the files and extract the text"""
    text_runs = []
    if path.endswith(".pdf"):
        loader = PyPDFLoader(path)
        pages = loader.load()
        print(pages)
        return pages
    elif path.endswith(".pptx"):
        prs = Presentation(path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
                        print(text_runs)
        return text_runs
    elif path.endswith('.jpg') or path.endswith('.jpeg') or path.endswith('.png'):
            image = cv2.imread(path)
            if image is not None:
                text = pytesseract.image_to_string(image)
                print(text)
            return text
    else:
        print("Incorrect file format, only supported file formats are pdf, jpg, jpeg, pptx, png.")

# text = load_file("data/Bhavya Kajani-1.pdf")
# print("FIle PRINTED ABOVE")

# text = normalize_text(load_file("path/to/resume.pdf"))
# chunks = get_chunks(text)


def get_chunks(text):
    """Split the extracted text into chunks"""
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size = 1500,
        chunk_overlap = 200,
        length_function = len,
        separators = ["\n\n", "\n", " "],
    )
    chunks_without_split = text # to compare with the split chunks
    chunks = text_splitter.split_text(text)
    return chunks

def normalize_text(raw_output) -> str:
    if isinstance(raw_output, list):
        if all(hasattr(p, 'page_content') for p in raw_output):
            return "\n".join([p.page_content for p in raw_output])
        return "\n".join(raw_output)
    return str(raw_output)


